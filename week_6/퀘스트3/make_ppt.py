# -*- coding: utf-8 -*-
"""
☕ Day On(데이온) — 고객의 소리(VOC) 분석 PPT 생성
실행: python make_ppt.py
입력: 09_고객의소리_리뷰원문.csv  (같은 폴더)
출력: 고객의소리_분석보고서.pptx

차트는 PowerPoint 네이티브 차트로 넣는다(이미지가 아니라 편집 가능).
"""
import csv, os
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "09_고객의소리_리뷰원문.csv")
OUT = os.path.join(HERE, "고객의소리_분석보고서.pptx")

# ── 브랜드 컬러 (내추럴·자연주의) ──
GREEN = RGBColor(0x2F, 0x5D, 0x50)   # 딥그린 (메인)
NUT   = RGBColor(0x8C, 0x62, 0x39)   # 너트브라운 (포인트)
BEIGE = RGBColor(0xEF, 0xE8, 0xDA)   # 베이지 (배경 박스)
GRAY  = RGBColor(0x6B, 0x6B, 0x6B)
DARK  = RGBColor(0x2B, 0x2B, 0x2B)
RED   = RGBColor(0xC0, 0x50, 0x4D)   # 불만
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "맑은 고딕"

# ── 테마 분류 규칙 (리뷰 1건이 여러 테마에 속할 수 있음) ──
THEMES = [
    ("시그니처·맛(견과)", ["너츠", "피칸", "고소", "시그니처", "맛있", "진하", "달지 않"]),
    ("공간·분위기",       ["식물", "인테리어", "원목", "조용", "편안", "아늑", "쉬어", "한 숨", "햇빛", "골목", "정원", "음악"]),
    ("재방문·단골",       ["재방문", "단골", "또 ", "자주", "매일", "루틴", "다섯 번째", "주 2회"]),
    ("서비스·사장님",     ["사장님", "친절", "응대", "눈치", "설명"]),
    ("가격·가성비",       ["비싸", "비쌈", "가격", "값은", "착해", "대비"]),
    ("좌석·대기",         ["자리", "좌석", "대기", "웨이팅", "돌아왔", "일찍"]),
    ("주차·접근성",       ["주차", "차 가져"]),
    ("테이크아웃·포장",   ["포장", "테이크아웃", "사가", "사갑니다"]),
]


def load_reviews():
    with open(SRC, encoding="utf-8-sig") as f:
        return list(csv.DictReader(f))


def tag_themes(text):
    hits = []
    for name, kws in THEMES:
        if any(k in text for k in kws):
            hits.append(name)
    return hits


# ─────────────────────────────────────────────
# 슬라이드 헬퍼
# ─────────────────────────────────────────────
def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃


def textbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, font=FONT, spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if spacing:
            p.space_after = Pt(spacing)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def slide_header(slide, title, sub=None):
    """상단 제목 + 좌측 그린 바"""
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(0.42), Inches(0.09), Inches(0.52))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background(); bar.shadow.inherit = False
    textbox(slide, 0.75, 0.35, 9.0, 0.6, title, size=26, bold=True, color=DARK)
    if sub:
        textbox(slide, 0.75, 0.95, 11.5, 0.4, sub, size=13, color=GRAY)


def style_chart(chart, size=11):
    chart.font.size = Pt(size)
    chart.font.name = FONT
    chart.font.color.rgb = DARK


# ─────────────────────────────────────────────
def build():
    rows = load_reviews()
    total = len(rows)

    ratings = Counter(int(r["평점"]) for r in rows)
    avg = sum(int(r["평점"]) for r in rows) / total

    # 채널별
    ch_cnt = Counter(r["채널"] for r in rows)
    ch_avg = {c: sum(int(r["평점"]) for r in rows if r["채널"] == c) / ch_cnt[c] for c in ch_cnt}
    ch_order = sorted(ch_cnt, key=lambda c: -ch_cnt[c])

    # 테마별 (긍정=4~5점, 불만=3점 이하)
    theme_pos, theme_neg = Counter(), Counter()
    for r in rows:
        pos = int(r["평점"]) >= 4
        for t in tag_themes(r["리뷰내용"]):
            (theme_pos if pos else theme_neg)[t] += 1
    theme_names = [t for t, _ in THEMES]

    # 메뉴별
    menu = defaultdict(list)
    for r in rows:
        m = r["언급메뉴"]
        if m and m != "(메뉴 언급 없음)":
            menu[m].append(int(r["평점"]))
    menu_top = sorted(menu.items(), key=lambda kv: (-len(kv[1]), -sum(kv[1]) / len(kv[1])))[:8]

    # 낮은 평점부터 = 심각한 것부터
    complaints = sorted((r for r in rows if int(r["평점"]) <= 3), key=lambda r: int(r["평점"]))

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9

    # ═══ 1. 표지 ═══
    s = add_slide(prs)
    rect(s, 0, 0, 13.333, 7.5, GREEN)
    textbox(s, 1.0, 2.3, 11, 0.6, "☕  Day On (데이온)", size=22, color=BEIGE)
    textbox(s, 1.0, 2.95, 11, 1.1, "고객의 소리 분석 보고서", size=44, bold=True, color=WHITE)
    textbox(s, 1.0, 4.15, 11, 0.5,
            "리뷰 60건 · 4개 채널 · 2026-04-14 ~ 2026-07-12 (90일)", size=16, color=BEIGE)
    rect(s, 1.0, 4.95, 3.0, 0.06, NUT)
    textbox(s, 1.0, 5.3, 11, 0.4, "별점 분포 · 테마별 집계 · 개선 현황", size=14, color=BEIGE)

    # ═══ 2. 한눈에 보기 (KPI) ═══
    s = add_slide(prs)
    slide_header(s, "한눈에 보기", "리뷰 60건 요약")
    pos_cnt = ratings[5] + ratings[4]
    kpis = [
        ("총 리뷰", f"{total}건", "4개 채널"),
        ("평균 평점", f"{avg:.2f}점", "5점 만점"),
        ("긍정 비율", f"{100*pos_cnt/total:.0f}%", f"4점 이상 {pos_cnt}건"),
        ("불만 건수", f"{len(complaints)}건", "3점 이하"),
    ]
    for i, (label, value, note) in enumerate(kpis):
        x = 0.75 + i * 3.1
        rect(s, x, 1.85, 2.85, 1.75, BEIGE)
        textbox(s, x + 0.25, 2.05, 2.4, 0.3, label, size=13, color=GRAY)
        textbox(s, x + 0.25, 2.4, 2.4, 0.7, value, size=32, bold=True, color=GREEN)
        textbox(s, x + 0.25, 3.12, 2.4, 0.3, note, size=11, color=GRAY)

    rect(s, 0.75, 4.05, 11.8, 2.5, WHITE, line=BEIGE)
    textbox(s, 1.05, 4.25, 11.2, 0.35, "핵심 메시지", size=15, bold=True, color=NUT)
    msg = ("1.  시그니처는 강하다 — 너츠라떼 8회 언급 평균 5.00점, 불만 0건. 견과 테마가 실제로 먹히고 있다.\n"
           "2.  불만은 메뉴가 아니라 '공간'에 있다 — 2점 리뷰 2건은 전부 좌석 부족·주차 문제.\n"
           "3.  일반 메뉴는 약하다 — 카페라떼 3.00점, 자몽에이드 3.50점. \"시그니처를 시키는 게 맞다\"는 리뷰까지 나왔다.")
    textbox(s, 1.05, 4.7, 11.2, 1.6, msg, size=14, color=DARK, spacing=10)

    # ═══ 3. 별점 분포 ═══
    s = add_slide(prs)
    slide_header(s, "별점 분포", f"평균 {avg:.2f}점 · 4점 이상이 {100*pos_cnt/total:.0f}%를 차지")
    cd = CategoryChartData()
    cd.categories = ["5점", "4점", "3점", "2점", "1점"]
    cd.add_series("리뷰 수", tuple(ratings.get(k, 0) for k in [5, 4, 3, 2, 1]))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(0.75), Inches(1.6), Inches(7.4), Inches(5.2), cd)
    ch = gf.chart
    ch.has_legend = False
    style_chart(ch, 12)
    pl = ch.plots[0]
    pl.has_data_labels = True
    pl.data_labels.font.size = Pt(13)
    pl.data_labels.font.bold = True
    pl.data_labels.font.name = FONT
    pl.gap_width = 60
    # 5·4점=그린, 3점=너트, 2·1점=레드
    pts = pl.series[0].points
    for i, c in enumerate([GREEN, GREEN, NUT, RED, RED]):
        pts[i].format.fill.solid()
        pts[i].format.fill.fore_color.rgb = c

    rect(s, 8.5, 1.6, 4.05, 5.2, BEIGE)
    textbox(s, 8.8, 1.85, 3.5, 0.35, "읽는 법", size=14, bold=True, color=NUT)
    txt = (f"· 5점 {ratings[5]}건 / 4점 {ratings[4]}건\n"
           f"· 3점 {ratings[3]}건 / 2점 {ratings[2]}건\n\n"
           "평점 자체는 건강하다.\n\n"
           "다만 2점 2건이 모두\n'좌석·주차' 문제다.\n"
           "맛이 아니라 공간 때문에\n손님을 잃고 있다.\n\n"
           "→ 매출 직접 손실")
    textbox(s, 8.8, 2.35, 3.5, 4.2, txt, size=13, color=DARK, spacing=4)

    # ═══ 4. 채널별 평점 ═══
    s = add_slide(prs)
    slide_header(s, "채널별 반응", "네이버는 후하고 카카오맵은 짜다 — 불만은 카카오맵에 모인다")
    cd = CategoryChartData()
    cd.categories = ch_order
    cd.add_series("평균 평점", tuple(round(ch_avg[c], 2) for c in ch_order))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(0.75), Inches(1.75), Inches(6.1), Inches(5.0), cd)
    ch = gf.chart
    ch.has_legend = False
    style_chart(ch, 12)
    ch.value_axis.maximum_scale = 5.0
    ch.value_axis.minimum_scale = 3.0
    pl = ch.plots[0]
    pl.has_data_labels = True
    pl.data_labels.font.size = Pt(12)
    pl.data_labels.font.bold = True
    pl.data_labels.font.name = FONT
    pl.gap_width = 80
    pl.series[0].format.fill.solid()
    pl.series[0].format.fill.fore_color.rgb = GREEN

    cd2 = CategoryChartData()
    cd2.categories = ch_order
    cd2.add_series("리뷰 수", tuple(ch_cnt[c] for c in ch_order))
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                             Inches(7.1), Inches(1.75), Inches(5.45), Inches(5.0), cd2)
    ch2 = gf2.chart
    ch2.has_legend = False
    style_chart(ch2, 12)
    pl2 = ch2.plots[0]
    pl2.has_data_labels = True
    pl2.data_labels.font.size = Pt(12)
    pl2.data_labels.font.name = FONT
    pl2.gap_width = 60
    pl2.series[0].format.fill.solid()
    pl2.series[0].format.fill.fore_color.rgb = NUT
    textbox(s, 0.75, 6.85, 5.5, 0.4, "◀ 채널별 평균 평점 (3~5점 구간)", size=11, color=GRAY)
    textbox(s, 7.1, 6.85, 5.5, 0.4, "◀ 채널별 리뷰 수", size=11, color=GRAY)

    # ═══ 5. 테마별 집계 ═══
    s = add_slide(prs)
    slide_header(s, "테마별 집계", "리뷰 내용을 8개 테마로 분류 (한 건이 여러 테마에 중복 집계)")
    # 테마명이 길어 세로축(가로막대)에 두어야 라벨이 잘리지 않는다. 언급 많은 순으로 정렬.
    ordered = sorted(theme_names, key=lambda t: theme_pos.get(t, 0) + theme_neg.get(t, 0))
    cd = CategoryChartData()
    cd.categories = ordered
    cd.add_series("긍정 (4~5점)", tuple(theme_pos.get(t, 0) for t in ordered))
    cd.add_series("불만 (3점 이하)", tuple(theme_neg.get(t, 0) for t in ordered))
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED,
                            Inches(0.75), Inches(1.75), Inches(8.3), Inches(4.6), cd)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.TOP
    ch.legend.include_in_layout = False
    style_chart(ch, 11)
    pl = ch.plots[0]
    pl.has_data_labels = True
    pl.data_labels.font.size = Pt(10)
    pl.data_labels.font.name = FONT
    pl.gap_width = 50
    pl.series[0].format.fill.solid(); pl.series[0].format.fill.fore_color.rgb = GREEN
    pl.series[1].format.fill.solid(); pl.series[1].format.fill.fore_color.rgb = RED

    rect(s, 9.35, 1.75, 3.2, 4.6, BEIGE)
    textbox(s, 9.6, 1.98, 2.8, 0.35, "해석", size=14, bold=True, color=NUT)
    t_txt = ("칭찬은\n'맛·공간'에 몰리고,\n\n불만은\n'좌석·주차·가격'에\n몰린다.\n\n"
             "즉 제품은 이겼고\n공간 운영이 발목을\n잡고 있다.")
    textbox(s, 9.6, 2.45, 2.8, 3.7, t_txt, size=13, color=DARK, spacing=3)
    textbox(s, 0.75, 6.5, 11.8, 0.4,
            "※ 키워드 기반 자동 분류 — 리뷰 1건이 여러 테마에 걸칠 수 있어 합계는 60을 넘는다.",
            size=10, color=GRAY)

    # ═══ 6. 메뉴별 리뷰 반응 ═══
    s = add_slide(prs)
    slide_header(s, "메뉴별 리뷰 반응", "시그니처는 5.00점 만점 — 일반 메뉴가 평균을 깎는다")
    cd = CategoryChartData()
    cd.categories = [m for m, _ in menu_top][::-1]
    cd.add_series("평균 평점", tuple(round(sum(v) / len(v), 2) for _, v in menu_top)[::-1])
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(0.75), Inches(1.75), Inches(8.3), Inches(4.9), cd)
    ch = gf.chart
    ch.has_legend = False
    style_chart(ch, 12)
    ch.value_axis.maximum_scale = 5.0
    ch.value_axis.minimum_scale = 2.5
    pl = ch.plots[0]
    pl.has_data_labels = True
    pl.data_labels.font.size = Pt(12)
    pl.data_labels.font.bold = True
    pl.data_labels.font.name = FONT
    pl.gap_width = 50
    pts = pl.series[0].points
    for i, (m, v) in enumerate(reversed(menu_top)):
        pts[i].format.fill.solid()
        avg_m = sum(v) / len(v)
        pts[i].format.fill.fore_color.rgb = GREEN if avg_m >= 4.5 else (NUT if avg_m >= 4 else RED)

    rect(s, 9.35, 1.75, 3.2, 4.9, BEIGE)
    textbox(s, 9.6, 1.98, 2.8, 0.35, "손님의 말", size=14, bold=True, color=NUT)
    quote = ('"라떼는 그냥 그랬어요.\n시그니처 메뉴를\n시키는 게 맞는 것\n같습니다."\n\n'
             "— 카카오맵 3점\n\n\n"
             "너츠라떼 5.00점\n간판세트 5.00점\n───────────\n"
             "카페라떼 3.00점\n자몽에이드 3.50점")
    textbox(s, 9.6, 2.45, 2.8, 4.0, quote, size=12, color=DARK, spacing=2)

    # ═══ 7. 불만 원문 (고객의 소리) ═══
    s = add_slide(prs)
    slide_header(s, "불만 사항 원문", f"3점 이하 {len(complaints)}건 — 손님이 실제로 쓴 말")
    rows_n = len(complaints) + 1
    tbl_shape = s.shapes.add_table(rows_n, 4, Inches(0.75), Inches(1.7),
                                   Inches(11.8), Inches(0.5 + 0.45 * len(complaints)))
    tbl = tbl_shape.table
    for i, w in enumerate([Inches(1.0), Inches(1.4), Inches(2.3), Inches(7.1)]):
        tbl.columns[i].width = w
    heads = ["평점", "채널", "언급 메뉴", "리뷰 내용"]
    for j, h in enumerate(heads):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = GREEN
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True
                r.font.color.rgb = WHITE; r.font.name = FONT
    for i, cm in enumerate(complaints, start=1):
        vals = [f"{cm['평점']}점", cm["채널"], cm["언급메뉴"], cm["리뷰내용"]]
        low = int(cm["평점"]) <= 2
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0xFA, 0xEC, 0xEB) if low else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = FONT
                    r.font.color.rgb = RED if (low and j == 0) else DARK
                    r.font.bold = (j == 0)

    # ═══ 8. 개선 현황 ═══
    s = add_slide(prs)
    slide_header(s, "개선 현황", "고객의 소리 → 개선 과제 (현재 전부 '계획' 단계 — 아직 실행 전)")
    heads = ["고객의 소리", "개선 과제", "상태", "근거 데이터"]
    items = [
        ("주말에 자리가 없어\n그냥 돌아왔다 (2점)",
         "만석 시 테이크아웃 전환 유도\n(피칸 토핑 무료, 가격 동일)",
         "계획",
         "주말 테이크아웃 28.5%\nvs 평일 42% → 전환 여지"),
        ("주차가 아예 안 된다 (2점)",
         "해결 불가 — 골목 상권 특성\n대신 '도보 5분' 안내 강화",
         "보류",
         "상권 구조상 개선 난이도 높음"),
        ("샌드위치 8,500원은\n비싸다 (3점)",
         "가격 인하 대신 런치 세트로 전환\n(샌드위치+아메리카노 10,500원)",
         "계획",
         "햄치즈 최근 30일 -5.7% 역성장"),
        ("라떼는 그냥 그랬다\n시그니처를 시켜라 (3점)",
         "비시그니처 메뉴 정리\n(제철 과일 티 단종 검토)",
         "계획",
         "과일 티 판매 최하위(70잔/30일)"),
        ("(리뷰엔 없지만 데이터로 확인)\n더우면 시그니처가 안 팔림",
         "신메뉴 '피칸 크림 콜드브루'\n6,500원 · 마진율 66%",
         "제안",
         "폭염일 매출 -18.3%\n간판세트 -31.9%"),
    ]
    tbl = s.shapes.add_table(len(items) + 1, 4, Inches(0.75), Inches(1.75),
                             Inches(11.8), Inches(4.9)).table
    for i, w in enumerate([Inches(3.0), Inches(3.5), Inches(1.1), Inches(4.2)]):
        tbl.columns[i].width = w
    for j, h in enumerate(heads):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = GREEN
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(13); r.font.bold = True
                r.font.color.rgb = WHITE; r.font.name = FONT
    status_color = {"계획": NUT, "보류": GRAY, "제안": GREEN}
    for i, row in enumerate(items, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.text = v
            c.fill.solid()
            c.fill.fore_color.rgb = BEIGE if i % 2 == 0 else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = FONT
                    r.font.bold = (j == 2)
                    r.font.color.rgb = status_color.get(v, DARK) if j == 2 else DARK
    textbox(s, 0.75, 6.8, 11.8, 0.4,
            "※ 상태는 모두 실행 전 단계다. '개선 완료' 항목은 없다 — 실제 조치 후 갱신 필요.",
            size=10, color=RED)

    # ═══ 9. 결론 ═══
    s = add_slide(prs)
    rect(s, 0, 0, 13.333, 7.5, GREEN)
    textbox(s, 1.0, 0.9, 11, 0.8, "결론", size=34, bold=True, color=WHITE)
    rect(s, 1.0, 1.85, 2.2, 0.05, NUT)

    concl = [
        ("제품은 이미 이겼다",
         "너츠라떼·간판세트 5.00점 만점, 불만 0건. 견과 테마 전략은 검증됐다. 건드리지 말 것."),
        ("발목을 잡는 건 공간이다",
         "2점 리뷰 2건 전부 좌석·주차. 맛이 아니라 자리가 없어서 손님을 돌려보내고 있다."),
        ("일반 메뉴가 평균을 깎는다",
         "카페라떼 3.00점, 자몽에이드 3.50점. 시그니처로 흡수하거나 정리해야 한다."),
    ]
    for i, (h, b) in enumerate(concl):
        y = 2.35 + i * 1.35
        rect(s, 1.0, y, 11.3, 1.1, RGBColor(0x3B, 0x6E, 0x5F))
        textbox(s, 1.35, y + 0.13, 10.6, 0.4, f"{i+1}.  {h}", size=18, bold=True, color=WHITE)
        textbox(s, 1.35, y + 0.58, 10.6, 0.4, b, size=13, color=BEIGE)

    textbox(s, 1.0, 6.6, 11.3, 0.4,
            "다음 액션 — ① 주말 테이크아웃 전환 유도  ② 피칸 크림 콜드브루 출시  ③ 피칸 재고 선발주",
            size=13, bold=True, color=BEIGE)

    prs.save(OUT)
    print(f"OK: {OUT}")
    print(f"   슬라이드 {len(prs.slides.__iter__.__self__._sldIdLst)}장 / 리뷰 {total}건 분석")


if __name__ == "__main__":
    build()

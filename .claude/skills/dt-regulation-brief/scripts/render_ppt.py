# -*- coding: utf-8 -*-
"""
render_ppt.py - slides.json 스펙을 임원 보고용 PPT(.pptx)로 렌더링

실행:
    py -3 render_ppt.py <slides.json> <출력.pptx>

슬라이드 디자이너 에이전트는 slides.json 까지만 쓰고, 실제 그리기는 이 스크립트가 한다.
그래야 매번 같은 브랜드·같은 품질로 나온다. 스펙은 references/slides-spec.md 참조.

지원 타입: bullets / table / chart / timeline / action  (+ 표지·마무리 자동 생성)
차트는 이미지가 아니라 PowerPoint 네이티브 차트라 발표 중에도 편집된다.
python-pptx 필요(없으면 자동 설치 시도).
"""
import json
import subprocess
import sys


def ensure_pptx():
    try:
        import pptx  # noqa
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "--quiet", "python-pptx"], check=True)


ensure_pptx()

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 브랜드 (금융기관 보고서: 네이비 + 시그널 컬러) ──
NAVY = RGBColor(0x1B, 0x33, 0x5F)   # 메인
BLUE = RGBColor(0x2E, 0x6B, 0xB8)   # 포인트
SKY = RGBColor(0xE8, 0xEF, 0xF7)    # 배경 박스
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x22, 0x26, 0x2B)
RED = RGBColor(0xC0, 0x39, 0x3B)    # 리스크·기한 강조
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 본문 영역
MARGIN = Inches(0.9)
BODY_TOP = Inches(1.85)
BODY_W = SLIDE_W - MARGIN * 2
BODY_H = SLIDE_H - BODY_TOP - Inches(0.8)

CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


# ─────────────────────────────── 공통 그리기 도우미

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 완전 빈 레이아웃


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def style(p, text, size, color=DARK, bold=False):
    """문단에 스타일 적용된 run 을 붙인다. (python-pptx 의 add_run() 은 인자를 받지 않는다)"""
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return run


def rect(slide, left, top, width, height, fill):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def slide_header(slide, title, index=None):
    """모든 본문 슬라이드 상단: 네이비 바 + 제목."""
    bar = slide.shapes.add_shape(1, MARGIN, Inches(0.75), Inches(0.09), Inches(0.5))  # 1 = RECTANGLE
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    tf = textbox(slide, MARGIN + Inches(0.25), Inches(0.66), BODY_W - Inches(1.0), Inches(0.75))
    p = tf.paragraphs[0]
    style(p, title, 26, NAVY, bold=True)

    if index is not None:
        tf2 = textbox(slide, SLIDE_W - MARGIN - Inches(0.7), Inches(6.75), Inches(0.6), Inches(0.35))
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        style(p2, str(index), 11, GRAY)


# ─────────────────────────────── 슬라이드 타입별

def render_cover(prs, spec):
    slide = blank(prs)
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(3.3))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False

    tf = textbox(slide, MARGIN, Inches(1.25), BODY_W, Inches(1.6))
    p = tf.paragraphs[0]
    style(p, spec.get("title", "제목 없음"), 40, WHITE, bold=True)
    if spec.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.space_before = Pt(10)
        style(p2, spec["subtitle"], 18, RGBColor(0xB8, 0xC8, 0xE0))

    meta = " · ".join(x for x in [spec.get("author", ""), spec.get("date", "")] if x)
    if meta:
        tf3 = textbox(slide, MARGIN, Inches(3.75), BODY_W, Inches(0.5))
        style(tf3.paragraphs[0], meta, 14, GRAY)


def render_bullets(slide, spec):
    tf = textbox(slide, MARGIN, BODY_TOP, BODY_W, BODY_H)
    first = True
    for b in spec.get("bullets", []):
        text = b if isinstance(b, str) else b.get("text", "")
        subs = [] if isinstance(b, str) else b.get("sub", [])

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        style(p, "• ", 20, BLUE, bold=True)
        style(p, text, 20, DARK)

        for s in subs:
            sp = tf.add_paragraph()
            sp.level = 1
            sp.space_after = Pt(6)
            style(sp, "– " + s, 15, GRAY)

    if spec.get("note"):
        slide.notes_slide.notes_text_frame.text = spec["note"]


def render_table(slide, spec):
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers or not rows:
        return
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = min(Inches(0.55) * n_rows, BODY_H)
    shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, BODY_TOP, BODY_W, height)
    table = shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        style(p, str(h), 14, WHITE, bold=True)

    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else SKY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            value = str(row[c]) if c < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            style(p, value, 13, DARK)


def render_chart(slide, spec):
    kind = CHART_TYPES.get(spec.get("chart", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    data = CategoryChartData()
    data.categories = spec.get("categories", [])
    for s in spec.get("series", []):
        data.add_series(s.get("name", ""), s.get("values", []))

    has_insight = bool(spec.get("insight"))
    chart_h = BODY_H - (Inches(0.95) if has_insight else Inches(0))
    frame = slide.shapes.add_chart(kind, MARGIN, BODY_TOP, BODY_W, chart_h, data)
    chart = frame.chart
    chart.font.name = FONT
    chart.font.size = Pt(13)

    series_count = len(spec.get("series", []))
    if series_count > 1 or kind == XL_CHART_TYPE.PIE:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.name = FONT

    if kind != XL_CHART_TYPE.PIE:
        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = BLUE if i == 0 else NAVY

    if has_insight:
        top = BODY_TOP + chart_h + Inches(0.15)
        box = rect(slide, MARGIN, top, BODY_W, Inches(0.7), SKY)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        style(p, "→ " + spec["insight"], 15, NAVY, bold=True)


def render_timeline(slide, spec):
    items = spec.get("items", [])[:5]
    if not items:
        return
    n = len(items)
    gap = Inches(0.25)
    card_w = int((BODY_W - gap * (n - 1)) / n)
    top = BODY_TOP + Inches(0.5)

    # 연결선
    line = slide.shapes.add_shape(1, MARGIN, top + Inches(0.62), BODY_W, Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = SKY
    line.line.fill.background()
    line.shadow.inherit = False

    for i, it in enumerate(items):
        left = MARGIN + (card_w + gap) * i
        dot = slide.shapes.add_shape(9, left + int(card_w / 2) - Inches(0.11), top + Inches(0.5),
                                     Inches(0.22), Inches(0.22))  # 9 = OVAL
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        dot.shadow.inherit = False

        tf = textbox(slide, left, top, card_w, Inches(0.5))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style(p, str(it.get("when", "")), 16, NAVY, bold=True)

        card = rect(slide, left, top + Inches(1.0), card_w, Inches(1.3), SKY)
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        style(cp, str(it.get("what", "")), 14, DARK)


def render_action(slide, spec):
    actions = spec.get("actions", [])[:4]
    top = BODY_TOP
    for i, a in enumerate(actions):
        box = rect(slide, MARGIN, top, BODY_W, Inches(0.95), SKY)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        style(p, f"{i + 1}. ", 18, RED, bold=True)
        style(p, str(a.get("what", "")), 18, DARK, bold=True)

        meta = " / ".join(x for x in [str(a.get("who", "")), str(a.get("when", ""))] if x)
        if meta:
            p2 = tf.add_paragraph()
            style(p2, "   담당·기한: " + meta, 13, GRAY)

        top += Inches(1.15)


def render_closing(prs):
    slide = blank(prs)
    band = slide.shapes.add_shape(1, 0, Inches(2.9), SLIDE_W, Inches(1.7))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tf = textbox(slide, MARGIN, Inches(3.25), BODY_W, Inches(1.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style(p, "감사합니다", 30, WHITE, bold=True)


RENDERERS = {
    "bullets": render_bullets,
    "table": render_table,
    "chart": render_chart,
    "timeline": render_timeline,
    "action": render_action,
}


def build(spec, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    render_cover(prs, spec)

    slides = spec.get("slides", [])
    for i, s in enumerate(slides, start=2):
        kind = s.get("type")
        fn = RENDERERS.get(kind)
        if not fn:
            raise ValueError(
                f"{i}번째 슬라이드의 type '{kind}' 은(는) 지원하지 않습니다. "
                f"지원 타입: {', '.join(RENDERERS)}"
            )
        slide = blank(prs)
        slide_header(slide, s.get("title", ""), index=i)
        fn(slide, s)

    render_closing(prs)
    prs.save(out_path)
    return len(slides) + 2


def main():
    if len(sys.argv) != 3:
        print("사용법: py -3 render_ppt.py <slides.json> <출력.pptx>")
        sys.exit(1)

    src, out = sys.argv[1], sys.argv[2]
    with open(src, encoding="utf-8-sig") as f:
        spec = json.load(f)

    total = build(spec, out)
    print(f"[OK] {out}  (총 {total}장: 표지 + 본문 {total - 2} + 마무리)")


if __name__ == "__main__":
    main()
